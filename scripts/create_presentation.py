#!/usr/bin/env python3
"""
Script pour créer une présentation PowerPoint animée
E-Commerce Symfony Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Fonction pour créer RGBColor compatible
def rgb_color(r, g, b):
    """Crée une couleur RGB"""
    return RGBColor(r, g, b)

# Couleurs du thème
COLORS = {
    'primary': (41, 128, 185),      # Bleu
    'secondary': (52, 73, 94),       # Gris foncé
    'accent': (46, 204, 113),        # Vert
    'warning': (241, 196, 15),       # Jaune
    'danger': (231, 76, 60),         # Rouge
    'light': (236, 240, 241),        # Gris clair
    'dark': (44, 62, 80),            # Noir
    'white': (255, 255, 255),        # Blanc
}

def add_animation(shape):
    """Ajoute une animation d'apparition à une forme"""
    # Animation XML pour effet Fade
    spTree = shape._element.getparent()
    timing = parse_xml(
        '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"/>'
        '</p:par></p:tnLst></p:timing>'
    )
    return timing

def set_shape_fill(shape, color):
    """Définit la couleur de remplissage d'une forme"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)

def add_title_slide(prs, title, subtitle=""):
    """Ajoute une slide de titre"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond dégradé (rectangle)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, COLORS['primary'])
    bg.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(*COLORS['light'])
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, title):
    """Ajoute une slide de section"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, COLORS['secondary'])
    bg.line.fill.background()
    
    # Numéro de section
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"SECTION {section_num}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(*COLORS['accent'])
    p.alignment = PP_ALIGN.CENTER
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, has_table=False, table_data=None):
    """Ajoute une slide de contenu"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre de titre
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    set_shape_fill(title_bar, COLORS['primary'])
    title_bar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    
    # Contenu
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*COLORS['dark'])
            p.space_after = Pt(12)
            p.level = 0
    
    # Tableau si fourni
    if has_table and table_data:
        rows = len(table_data)
        cols = len(table_data[0])
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3.5), Inches(9), Inches(2.5)).table
        
        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                if i == 0:
                    cell.text_frame.paragraphs[0].font.bold = True
                    set_cell_fill(cell, COLORS['primary'])
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*COLORS['white'])
    
    return slide

def set_cell_fill(cell, color):
    """Définit la couleur de fond d'une cellule de tableau"""
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*color)

def add_table_slide(prs, title, table_data, subtitle=""):
    """Ajoute une slide avec tableau"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre de titre
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    set_shape_fill(title_bar, COLORS['primary'])
    title_bar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    
    # Sous-titre
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*COLORS['secondary'])
    
    # Tableau
    top = Inches(2) if subtitle else Inches(1.5)
    rows = len(table_data)
    cols = len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.3), top, Inches(9.4), Inches(4)).table
    
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER
            
            if i == 0:
                para.font.bold = True
                set_cell_fill(cell, COLORS['primary'])
                para.font.color.rgb = RGBColor(*COLORS['white'])
            else:
                set_cell_fill(cell, COLORS['light'] if i % 2 == 0 else COLORS['white'])
    
    return slide

def add_comparison_slide(prs, title, before_items, after_items):
    """Ajoute une slide de comparaison Avant/Après"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre de titre
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    set_shape_fill(title_bar, COLORS['primary'])
    title_bar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    
    # Colonne Avant (rouge)
    before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(4.5))
    set_shape_fill(before_box, COLORS['danger'])
    
    before_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(4), Inches(0.5))
    tf = before_title.text_frame
    p = tf.paragraphs[0]
    p.text = "AVANT"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    before_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(4), Inches(3.5))
    tf = before_content.text_frame
    for i, item in enumerate(before_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"- {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.space_after = Pt(10)
    
    # Colonne Après (vert)
    after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(4.5))
    set_shape_fill(after_box, COLORS['accent'])
    
    after_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4), Inches(0.5))
    tf = after_title.text_frame
    p = tf.paragraphs[0]
    p.text = "APRES"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    after_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.3), Inches(4), Inches(3.5))
    tf = after_content.text_frame
    for i, item in enumerate(after_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"+ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.space_after = Pt(10)
    
    return slide

def add_metrics_slide(prs, title, metrics):
    """Ajoute une slide de métriques avec grands chiffres"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre de titre
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    set_shape_fill(title_bar, COLORS['primary'])
    title_bar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    
    # Métriques en grille
    cols = min(len(metrics), 3)
    box_width = 2.8
    start_x = (10 - (cols * box_width + (cols-1) * 0.3)) / 2
    
    for i, (value, label, color) in enumerate(metrics):
        col = i % cols
        row = i // cols
        x = start_x + col * (box_width + 0.3)
        y = 1.8 + row * 2.2
        
        # Boîte
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_width), Inches(2))
        set_shape_fill(box, color)
        
        # Valeur
        value_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(box_width), Inches(1))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.2), Inches(box_width), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*COLORS['white'])
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_conclusion_slide(prs, points):
    """Ajoute une slide de conclusion"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, COLORS['secondary'])
    bg.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Points Cles a Retenir"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    # Points
    for i, point in enumerate(points):
        y = 1.8 + i * 0.9
        point_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.8))
        set_shape_fill(point_box, COLORS['accent'] if i % 2 == 0 else COLORS['primary'])
        
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.2), Inches(8.6), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[OK] {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*COLORS['white'])
    
    return slide

def add_thank_you_slide(prs):
    """Ajoute une slide de remerciement"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, COLORS['primary'])
    bg.line.fill.background()
    
    # Merci
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Merci !"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*COLORS['white'])
    p.alignment = PP_ALIGN.CENTER
    
    # Questions
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions ?"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(*COLORS['light'])
    p.alignment = PP_ALIGN.CENTER
    
    # Info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Projet E-Commerce Symfony - Decembre 2025"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(*COLORS['light'])
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Crée la présentation complète"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ============================================
    # SLIDE 0: Titre
    # ============================================
    add_title_slide(prs, 
        "Plateforme E-Commerce Symfony",
        "Architecture Sécurisée & Haute Disponibilité"
    )
    
    # ============================================
    # SECTION 1: Contexte / Problématique
    # ============================================
    add_section_slide(prs, 1, "Contexte / Problématique")
    
    # Slide 1.1: Contexte
    add_content_slide(prs, "Slide 1.1 : Contexte du Projet", [
        "🎯 OBJECTIF PRINCIPAL",
        "Développer une plateforme e-commerce sécurisée, scalable et moderne",
        "",
        "📋 BESOINS IDENTIFIÉS",
        "• 🛒 Catalogue Produits : Gestion complète des produits et catégories",
        "• 👤 Gestion Utilisateurs : Inscription, connexion, rôles (Admin/Client)",
        "• 🔐 Sécurité Renforcée : Protection contre les cyberattaques",
        "• 💳 Commandes : Panier, checkout, historique",
        "• 📱 Accessibilité : Interface responsive (mobile/desktop)",
        "",
        "🏢 CONTEXTE TECHNIQUE",
        "• Stack : PHP 8.2 / Symfony 5.4 / MySQL 8.0",
        "• Infrastructure : Docker + Kubernetes",
        "• Environnement : Développement local → Production cloud"
    ])
    
    # Slide 1.2: Problématiques
    add_content_slide(prs, "Slide 1.2 : Problématiques Identifiées", [
        "⚠️ DÉFIS MAJEURS",
        "",
        "🔐 SÉCURITÉ",
        "• Authentification faible (mots de passe simples)",
        "• Attaques par force brute",
        "• Données sensibles exposées",
        "",
        "📈 SCALABILITÉ & DISPONIBILITÉ",
        "• Charge variable (pics de trafic Black Friday)",
        "• Single point of failure",
        "• Temps d'arrêt lors des maintenances",
        "",
        "📊 STATISTIQUES DU MARCHÉ",
        "• 43% des cyberattaques ciblent les PME",
        "• 60% ferment dans les 6 mois après une attaque",
        "• 94% des utilisateurs abandonnent un site non sécurisé"
    ])
    
    # ============================================
    # SECTION 2: Solutions Possibles
    # ============================================
    add_section_slide(prs, 2, "Solutions Possibles")
    
    # Slide 2.1: Solutions Sécurité
    add_table_slide(prs, "Slide 2.1 : Analyse des Solutions de Sécurité",
        [
            ["Solution", "Sécurité", "Complexité", "Coût", "Score"],
            ["Mot de passe simple", "⭐ (20%)", "Faible", "0€", "2/10"],
            ["MDP + Politique", "⭐⭐ (40%)", "Moyenne", "0€", "4/10"],
            ["2FA (TOTP) ✅", "⭐⭐⭐⭐⭐ (99.9%)", "Moyenne", "0€", "9/10"],
            ["Biométrique", "⭐⭐⭐⭐⭐ (99.9%)", "Haute", "€€€", "7/10"],
            ["SSO (OAuth2)", "⭐⭐⭐⭐ (95%)", "Haute", "€€", "6/10"],
        ],
        "Comparaison des méthodes d'authentification"
    )
    
    # Slide 2.2: Solutions Infrastructure
    add_table_slide(prs, "Slide 2.2 : Analyse des Solutions d'Infrastructure",
        [
            ["Architecture", "Disponibilité", "Scalabilité", "Coût/mois", "Score"],
            ["Serveur unique", "95%", "Faible", "~20€", "3/10"],
            ["Docker + K8s ✅", "99.5%", "Haute", "~50€", "8/10"],
            ["Multi-cluster + LB", "99.99%", "Très haute", "~150€", "9/10"],
            ["Cloud managé (AWS)", "99.99%", "Auto", "~200€+", "9/10"],
        ],
        "Comparaison des architectures"
    )
    
    # ============================================
    # SECTION 3: Déroulement / Réalisation
    # ============================================
    add_section_slide(prs, 3, "Déroulement / Réalisation")
    
    # Slide 3.1: Architecture Technique
    add_table_slide(prs, "Slide 3.1 : Architecture Technique Implémentée",
        [
            ["Couche", "Technologie", "Version"],
            ["Frontend", "Twig + Bootstrap", "5.x"],
            ["Backend", "Symfony", "5.4"],
            ["Langage", "PHP", "8.2"],
            ["Base de données", "MySQL", "8.0"],
            ["Conteneurisation", "Docker", "24.x"],
            ["Orchestration", "Kubernetes (Minikube)", "1.28"],
            ["2FA", "scheb/2fa-bundle + endroid/qr-code", "6.x"],
            ["SSL", "Let's Encrypt", "-"],
        ],
        "Stack technologique complète"
    )
    
    # Slide 3.2: Fonctionnalités
    add_content_slide(prs, "Slide 3.2 : Fonctionnalités Développées", [
        "👥 GESTION DES UTILISATEURS",
        "Inscription → Validation → Connexion → 2FA QR Code → OTP → Accès",
        "",
        "🛒 PROCESSUS D'ACHAT",
        "• ✅ Parcourir le catalogue",
        "• ✅ Rechercher des produits",
        "• ✅ Ajouter au panier",
        "• ✅ Authentification 2FA",
        "• ✅ Valider la commande",
        "• ✅ Paiement",
        "",
        "👨‍💼 FONCTIONNALITÉS ADMIN",
        "• CRUD Produits • Gestion Catégories • Gestion Utilisateurs",
        "• Tableau de bord • Rapports de vente",
        "",
        "🔐 SÉCURITÉ : 2FA TOTP • HTTPS • Politique MDP • Firewall • Sauvegardes"
    ])
    
    # ============================================
    # SECTION 4: Bilan
    # ============================================
    add_section_slide(prs, 4, "Bilan (Analyse des Résultats)")
    
    # Slide 4.1: Résultats
    add_metrics_slide(prs, "Slide 4.1 : Résultats Obtenus", [
        ("99.7%", "Disponibilité", COLORS['accent']),
        ("0.8s", "Temps de réponse", COLORS['primary']),
        ("100%", "Couverture 2FA", COLORS['accent']),
        ("100%", "HTTPS", COLORS['primary']),
        ("669%", "ROI", COLORS['warning']),
        ("4 pods", "Scalabilité", COLORS['accent']),
    ])
    
    # Slide 4.2: Avant/Après
    add_comparison_slide(prs, "Slide 4.2 : Amélioration de la Sécurité",
        [
            "Mot de passe simple",
            "HTTP non chiffré",
            "Serveur unique (95%)",
            "Backup manuel",
            "Risque: 15,000€/an"
        ],
        [
            "2FA TOTP (+99.9%)",
            "HTTPS Let's Encrypt",
            "Multi-pod K8s (99.5%)",
            "Backup automatique",
            "Risque: 15€/an"
        ]
    )
    
    # Slide Conclusion
    add_conclusion_slide(prs, [
        "Sécurité 2FA : Protection maximale des comptes utilisateurs",
        "Architecture K8s : Haute disponibilité et scalabilité",
        "Automatisation : Scripts de déploiement et rebuild",
        "ROI de 669% : Investissement rentabilisé",
        "Standards respectés : OWASP, 12-Factor App"
    ])
    
    # Slide Merci
    add_thank_you_slide(prs)
    
    # Sauvegarde
    output_path = os.path.join(os.path.dirname(__file__), 'presentation_ecommerce.pptx')
    prs.save(output_path)
    print(f"✅ Présentation créée : {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
